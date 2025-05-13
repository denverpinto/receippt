import os
import json
from pptx import Presentation 

# index created for directories -
slidesRootDir = "./slides"
noteTagLabel = "RECEIPPT-TAGS:"
templatesRootDir = "./templates"

# create updated index from scratch
updatedIndex = { "slides": [], "templates": [] }

# function to return all slides' text as html
def getSlidesTextAsHTML(path, verseName,tags):
	prs = Presentation(path)
	texts = []
	texts.append("<hr>")
	texts.append(f"<h4> VERSE: {verseName} <i>[{', '.join(tags)}]</i> </h4>")
	texts.append("<hr>")
	for slide_number, slide in enumerate(prs.slides):
		firstShape = True
		for shape in slide.shapes: 
			if hasattr(shape, "text"):
				if firstShape:
					if slide_number == 0:
						texts.append(f"<h4>{shape.text}</h4>")
					firstShape = False
				else:
					texts.append(f"{shape.text}")
	return "</br>".join(texts).replace("\n","</br>")

# function to return all slides' text
def getSlidesText(path, verseName, tags):
	prs = Presentation(path)
	texts = []
	texts.append("-----")
	for slide_number, slide in enumerate(prs.slides):
		firstShape = True
		for shape in slide.shapes: 
			if hasattr(shape, "text"):
				if firstShape:
					firstShape = False
				else:
					texts.append(f"{shape.text}")
	return "\n".join(texts)

# function to remove unknown slides and verses from templates TODO:
def scrubTemplate(template):
	knownSlides = [slideDetails["name"] for slideDetails in updatedIndex["slides"]]
	# scrub unknown slides
	for masspart in template["massparts"]:
		masspart["slides"] = [slideDetails for slideDetails in masspart["slides"] if slideDetails["name"] in knownSlides]
	# scrub unknown slide verses
	for masspart in template["massparts"]:
		for slide in masspart["slides"]:
			matchingIndexSlide = next((indexSlide for indexSlide in updatedIndex["slides"] if indexSlide["name"] == slide["name"]), None)
			if matchingIndexSlide:
				# slide["verses"] = [verseDetails for verseDetails in slide["verses"] if verseDetails["name"] in matchingIndexSlide["desiredVerses"]]
				slide["desiredVerses"] = [desiredVerse for desiredVerse in slide["desiredVerses"] if desiredVerse in matchingIndexSlide["desiredVerses"]]
				slide["verses"] = matchingIndexSlide["verses"]
				slide["path"] = matchingIndexSlide["path"]
				slide["tags"] = matchingIndexSlide["tags"]
	# scrub slides with empty desiredVerses
	for masspart in template["massparts"]:
		masspart["slides"] = [slideDetails for slideDetails in masspart["slides"] if len(slideDetails["desiredVerses"]) > 0]

	return template


# retrieve existing index if it exists
try:
	with open( './index.json','r') as f:
		currentIndex = json.load(f)
		print("Found Existing Index With {} Slides, {} Templates and {} Verses ".format(len(currentIndex["slides"]),len(currentIndex["templates"]),sum(len(slide["verses"]) for slide in currentIndex["slides"])))
except IOError:
	print("No Index Found, Creating Index From Scratch")

for hymnFolder in os.listdir(slidesRootDir):
	hymnPath = os.path.join(slidesRootDir, hymnFolder)
	if os.path.isdir(hymnPath):
		hymnEntry = {}
		hymnEntry["name"] = " ".join(" ".join(hymnFolder.split("_")).split()).upper()
		hymnEntry["path"] = hymnPath[2:]
		hymnEntry["tags"] = []
		hymnEntry["verses"] = []

		for verseFile in os.listdir(hymnPath):
			versePath = os.path.join(hymnPath, verseFile)
			verseEntry = {}
			verseEntry["name"] = " ".join(" ".join(verseFile.split("_")).split()).split(".pptx")[0].upper()
			verseEntry["path"] = versePath[2:]

			prs = Presentation(versePath)
			for slide_number, slide in enumerate(prs.slides):
				if slide_number == 0:
					noteText = slide.notes_slide.notes_text_frame.text
					if noteText.strip().upper().startswith(noteTagLabel):
						tags = noteText.strip().split(noteTagLabel)[1].strip().split(",")
						tags = [tag.strip().upper() for tag in tags]
						tags = [tag for tag in tags if tag]
						verseEntry["tags"] = tags
					else:
						verseEntry["tags"] = []
			verseEntry["html"] = getSlidesTextAsHTML(verseEntry["path"], verseEntry["name"], verseEntry["tags"])
			verseEntry["text"] = getSlidesText(verseEntry["path"], verseEntry["name"], verseEntry["tags"])

			hymnEntry["verses"].append(verseEntry)
			hymnEntry["tags"] = list(set(hymnEntry["tags"]) | set(verseEntry["tags"]))
			
		hymnEntry["desiredVerses"] = [verse["name"] for verse in hymnEntry["verses"]]
		updatedIndex["slides"].append(hymnEntry)		

# add templates 
for (root,dirs,files) in os.walk(templatesRootDir, topdown=False):
	if root == templatesRootDir : # top level directory containing template jsons
		for file in files:
			try:
				with open(root+"/"+file,'r') as f:
					currentTemplate = json.load(f)
					updatedIndex["templates"].append(scrubTemplate(currentTemplate))
			except IOError:
				print("{} couldn't be loaded".format(root+"/"+file))

# sort index.slides alphabetically
def sortFn(e):
  return e['name']

updatedIndex["slides"].sort(key=sortFn)

# display updated stats
verses = [verse for slide in updatedIndex["slides"] for verse in slide["verses"]]
print("Updated Index Has {} Slides, {} Templates and {} Verses \n".format(len(updatedIndex["slides"]),len(updatedIndex["templates"]),len(verses),sum(len(slide["verses"]) for slide in updatedIndex["slides"])))

# persist index to file
with open('./index.json', 'w') as f:
    json.dump(updatedIndex, f,indent = 2, sort_keys=True)  # encode dict into JSON
