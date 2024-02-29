import os
import json
from pptx import Presentation 

# index created for directory -
rootDir = "./slides"

# retrieve existing index if it exists
try:
	with open( './index.json','r') as f:
		currentIndex = json.load(f)
		print("Found Existing Index With {} Slides ".format(len(currentIndex)))
except IOError:
	currentIndex = {}
	print("No Index Found, Creating Index From Scratch")


# create updated index from scratch
updatedIndex = {}

for (root,dirs,files) in os.walk(rootDir, topdown=False):
	if root == rootDir : # top level directory containing folder tag names
		for file in files:
			updatedIndex[file] = {}
			updatedIndex[file]["name"] = file
			prs = prs = Presentation(root+"/"+file)
			texts = []
			for slide_number, slide in enumerate(prs.slides):
				texts.append(f"Slide {slide_number + 1}:")
				for shape in slide.shapes: 
					if hasattr(shape, "text"):
						texts.append(shape.text)
				if slide_number == 0:
					updatedIndex[file]["tags"] = slide.notes_slide.notes_text_frame.text
			updatedIndex[file]["text"] = "\n".join(texts)

# display updated stats
print("Updated Index Has {} Slides \n".format(len(updatedIndex)))


# persist index to file
with open('./index.json', 'w') as f:
    json.dump(updatedIndex, f,indent = 2, sort_keys=True)  # encode dict into JSON
